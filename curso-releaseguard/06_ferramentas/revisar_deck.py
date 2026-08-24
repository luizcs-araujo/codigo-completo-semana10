#!/usr/bin/env python3
"""Revisa o deck ReleaseGuard a partir do guia de aula slide a slide.

O script preserva os slides de evidência/código que já estavam corretos, substitui
blocos genéricos ou duplicados por conteúdo específico e injeta o guia completo
nas notas do apresentador.
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ACCENTS = {
    1: RGBColor(242, 140, 0),
    2: RGBColor(208, 32, 140),
    3: RGBColor(37, 70, 210),
}
BLACK = RGBColor(12, 12, 12)
DARK = RGBColor(28, 28, 30)
WHITE = RGBColor(250, 250, 250)
MUTED = RGBColor(190, 190, 194)


# Slides cujo conteúdo visível era genérico, repetido ou incompatível com o guia.
REBUILD = {
    4, 5, 7, 9, 13, 16, 17, 19, 21, 24, 25, 26, 27,
    29, 30, 33, 40, 41, 43, 45, 46, 47, 48, 49, 50, 52, 53, 54,
    56, 57, 59, 60, 63, 65, 69, 70, 71, 73, 74, 77, 78, 79, 80,
}


OVERRIDES = {
    25: {
        "primary": (
            "8 riscos, 8 grupos: carrinho stateful; boundary de estoque; checkout vazio; "
            "pagamento divergente; payment timeout; inventory 500; host/método proibido; plano Ollama."
        ),
        "engineering": "Entregar TestPlan válido, oracle justificado, HTTP real, report.json e uma limitação identificada.",
        "caution": "Criar em student_work/day1/<grupo>. Não alterar app/, schemas, policy, executor ou generate_plan.",
        "question": "Qual caso exige mais estado? Qual exige mais controle de policy?",
    },
    26: {
        "primary": "Oracle fundamentado 25% · plano mínimo 20% · execução HTTP real 20% · estado/placeholders 15%.",
        "engineering": "Policy/segurança 10% · explicação técnica 10%. “Funcionou” não basta: é preciso provar por que passou.",
        "caution": "A avaliação privilegia decisão técnica, evidência e rastreabilidade — não memorização nem adaptação do núcleo.",
        "question": "Como provar que o teste não passou por acidente?",
    },
    52: {
        "primary": (
            "8 variações: CTA deslocado; CTA ausente; conteúdo dinâmico; cor; viewport; "
            "fonte/rasterização; região crítica; desacordo VLM × threshold."
        ),
        "engineering": "Entregar baseline, current, diff, metrics.json, policy e justificativa de false positive/negative.",
        "caution": "Trabalhar em student_work/day2/<grupo>; não alterar compare.py, vlm_triage.py nem o estado do núcleo.",
        "question": "Qual caso tem maior risco de falso negativo?",
    },
    53: {
        "primary": "Policy alinhada ao risco 25% · captura reproduzível 20% · métricas bem interpretadas 20%.",
        "engineering": "FP/FN 15% · evidência visual 10% · explicação 10%. Score alto não compensa uma policy indefensável.",
        "caution": "Registrar viewport, browser, scale e cenário. Sem ambiente controlado, a comparação não é reproduzível.",
        "question": "Como o grupo provará que o ambiente foi controlado?",
    },
    79: {
        "primary": (
            "8 casos: payment_latency; payment_timeout; inventory_500; inventory_slow; "
            "ausência de mudanças; telemetria insuficiente; autonomia; release gate."
        ),
        "engineering": "Entregar sintoma → hipótese → tools/queries → evidência → revisão → confiança → HITL.",
        "caution": "Não começar pela causa. Não executar mutação. Evidência insuficiente é uma conclusão válida.",
        "question": "Qual grupo provavelmente terminará com menor confiança — e por quê?",
    },
    80: {
        "primary": "Contratos · instrumentação · identidade de serviço · tools limitadas · menor privilégio.",
        "engineering": "Auditoria de tool calls · avaliações · HITL · release policy · artefatos versionáveis.",
        "caution": "Checklist não substitui threat modeling, desenho de SLO nem validação específica do contexto produtivo.",
        "question": "Qual item vocês acrescentariam para o sistema real da empresa?",
    },
}


def clean_markdown(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def parse_guide(path: Path) -> dict[int, dict[str, str]]:
    raw = path.read_text(encoding="utf-8")
    matches = list(re.finditer(r"^## Slide (\d+) — (.+)$", raw, flags=re.MULTILINE))
    slides: dict[int, dict[str, str]] = {}
    labels = {
        # Metadados invisíveis no Markdown mantêm a reconstrução dos cards do
        # deck desacoplada da fala longa preparada para o professor.
        "primary": r"<!--\s*deck-primary:\s*(.+?)\s*-->",
        "engineering": r"<!--\s*deck-engineering:\s*(.+?)\s*-->",
        "caution": r"<!--\s*deck-caution:\s*(.+?)\s*-->",
        "question": r"<!--\s*deck-question:\s*(.+?)\s*-->",
    }
    for pos, match in enumerate(matches):
        number = int(match.group(1))
        end = matches[pos + 1].start() if pos + 1 < len(matches) else len(raw)
        section = raw[match.start():end].strip()
        fields = {"title": match.group(2).strip(), "section": section}
        for key, pattern in labels.items():
            found = re.search(pattern, section, flags=re.DOTALL)
            fields[key] = clean_markdown(found.group(1)) if found else ""
        slides[number] = fields
    if len(slides) != 81:
        raise ValueError(f"Guia incompleto: esperados 81 slides, encontrados {len(slides)}")
    return slides


def add_textbox(slide, x, y, w, h, text, size, color, *, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, fill, line, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(1.2)
    return card


def body_size(text: str, maximum: int = 24) -> int:
    if len(text) > 260:
        return min(maximum, 17)
    if len(text) > 190:
        return min(maximum, 19)
    if len(text) > 125:
        return min(maximum, 21)
    return maximum


def remove_content(slide):
    for shape in list(slide.shapes):
        if Inches(1.15) <= shape.top < Inches(6.35):
            shape._element.getparent().remove(shape._element)


def rebuild_slide(slide, day: int, data: dict[str, str]):
    remove_content(slide)
    accent = ACCENTS[day]

    add_card(slide, Inches(0.62), Inches(1.38), Inches(7.05), Inches(3.55), DARK, DARK)
    add_card(slide, Inches(0.62), Inches(1.38), Inches(0.08), Inches(3.55), accent, accent, radius=False)
    add_textbox(slide, Inches(0.88), Inches(1.62), Inches(6.45), Inches(0.34), "IDEIA CENTRAL", 12, accent, bold=True)
    add_textbox(
        slide, Inches(0.88), Inches(2.03), Inches(6.42), Inches(2.42),
        data["primary"], body_size(data["primary"], 26), WHITE, bold=True,
    )

    add_card(slide, Inches(7.98), Inches(1.38), Inches(4.72), Inches(1.78), WHITE, WHITE)
    add_textbox(slide, Inches(8.20), Inches(1.56), Inches(4.25), Inches(0.28), "DECISÃO DE ENGENHARIA", 11, accent, bold=True)
    add_textbox(
        slide, Inches(8.20), Inches(1.92), Inches(4.22), Inches(0.96),
        data["engineering"], body_size(data["engineering"], 17), BLACK,
    )

    add_card(slide, Inches(7.98), Inches(3.35), Inches(4.72), Inches(1.58), BLACK, accent)
    add_textbox(slide, Inches(8.20), Inches(3.54), Inches(4.25), Inches(0.28), "LIMITE / CUIDADO", 11, accent, bold=True)
    add_textbox(
        slide, Inches(8.20), Inches(3.88), Inches(4.22), Inches(0.82),
        data["caution"], body_size(data["caution"], 16), MUTED,
    )

    add_card(slide, Inches(0.62), Inches(5.22), Inches(12.08), Inches(0.74), accent, accent)
    add_textbox(
        slide, Inches(0.86), Inches(5.30), Inches(11.60), Inches(0.55),
        "PERGUNTA À TURMA  ·  " + data["question"], body_size(data["question"], 16), BLACK, bold=True,
    )


def add_combined_card(slide, x, y, w, h, label, body, fill, line, label_color, body_color, body_font):
    card = add_card(slide, x, y, w, h, fill, line)
    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.28)
    frame.margin_right = Inches(0.24)
    frame.margin_top = Inches(0.18)
    frame.margin_bottom = Inches(0.12)
    first = frame.paragraphs[0]
    first.space_after = Pt(12)
    label_run = first.add_run()
    label_run.text = label
    label_run.font.name = "Arial"
    label_run.font.size = Pt(11)
    label_run.font.bold = True
    label_run.font.color.rgb = label_color
    second = frame.add_paragraph()
    body_run = second.add_run()
    body_run.text = body
    body_run.font.name = "Arial"
    body_run.font.size = Pt(body_font)
    body_run.font.bold = body_font >= 20
    body_run.font.color.rgb = body_color
    return card


def rebuild_slide_combined(slide, day: int, data: dict[str, str]):
    """Variante robusta para páginas que o Keynote recorta ao importar shapes sobrepostos."""
    remove_content(slide)
    accent = ACCENTS[day]
    def text_card(x, y, w, h, label, body, fill, line, label_color, body_color, font_size):
        box = slide.shapes.add_textbox(x, y, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = line
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.28)
        frame.margin_right = Inches(0.24)
        frame.margin_top = Inches(0.18)
        first = frame.paragraphs[0]
        first.space_after = Pt(12)
        label_run = first.add_run()
        label_run.text = label
        label_run.font.name = "Arial"
        label_run.font.size = Pt(11)
        label_run.font.bold = True
        label_run.font.color.rgb = label_color
        second = frame.add_paragraph()
        body_run = second.add_run()
        body_run.text = body
        body_run.font.name = "Arial"
        body_run.font.size = Pt(font_size)
        body_run.font.bold = font_size >= 20
        body_run.font.color.rgb = body_color

    text_card(
        Inches(0.62), Inches(1.38), Inches(7.05), Inches(3.55),
        "IDEIA CENTRAL", data["primary"], DARK, accent, accent, WHITE, body_size(data["primary"], 25),
    )
    text_card(
        Inches(7.98), Inches(1.38), Inches(4.72), Inches(1.78),
        "DECISÃO DE ENGENHARIA", data["engineering"], WHITE, WHITE, accent, BLACK, body_size(data["engineering"], 17),
    )
    text_card(
        Inches(7.98), Inches(3.35), Inches(4.72), Inches(1.58),
        "LIMITE / CUIDADO", data["caution"], BLACK, accent, accent, MUTED, body_size(data["caution"], 16),
    )
    text_card(
        Inches(0.62), Inches(5.22), Inches(12.08), Inches(0.74),
        "", "PERGUNTA À TURMA  ·  " + data["question"], accent, accent, accent, BLACK,
        body_size(data["question"], 16),
    )


def remove_duplicate_page_number(slide, number: int):
    candidates = [
        shape for shape in slide.shapes
        if getattr(shape, "text", "").strip() == str(number) and shape.top > Inches(6.2)
    ]
    for shape in sorted(candidates, key=lambda item: item.left)[:-1]:
        shape._element.getparent().remove(shape._element)


def set_notes(slide, number: int, guide: dict[str, str], day: int):
    sources = {
        1: "00_guia_geral_3_dias.md; 01_dia1_implementacao.md; 06_guia_n8n_passo_a_passo.md; 01_dia1_exercicio_mentorado.md",
        2: "00_guia_geral_3_dias.md; 02_dia2_implementacao.md; 02_dia2_exercicio_mentorado.md",
        3: "00_guia_geral_3_dias.md; 03_dia3_implementacao.md; 03_dia3_teoria_conducao.md; 03_dia3_exercicio_mentorado.md",
    }
    # Os comentários guardam apenas textos curtos usados na reconstrução dos
    # cards visíveis e não devem aparecer nas notas do apresentador.
    note_section = re.sub(r"<!--.*?-->", "", guide["section"], flags=re.DOTALL).strip()
    note = (
        f"GUIA DO PROFESSOR — SLIDE {number:02d}\n\n"
        f"{note_section}\n\n"
        f"FONTES COMPLEMENTARES CONSIDERADAS\n{sources[day]}\n\n"
        "Princípio de condução: executar/observar antes de formalizar; separar proposta probabilística de validação, evidência e policy determinísticas."
    )
    slide.notes_slide.notes_text_frame.text = note


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("guide", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    guides = parse_guide(args.guide)
    presentation = Presentation(args.source)
    if len(presentation.slides) != 81:
        raise ValueError(f"Deck inesperado: {len(presentation.slides)} slides")

    # Corrige uma referência factual inexistente sem alterar o layout da capa.
    for shape in presentation.slides[0].shapes:
        if hasattr(shape, "text_frame") and "releaseguard_course_validated_2026-08-17.zip" in shape.text:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace(
                        "releaseguard_course_validated_2026-08-17.zip",
                        "nucleo_releaseguard_validado.zip",
                    )

    for number, slide in enumerate(presentation.slides, start=1):
        day = 1 if number <= 27 else 2 if number <= 54 else 3
        remove_duplicate_page_number(slide, number)
        set_notes(slide, number, guides[number], day)
        if number in REBUILD:
            content = {
                "primary": guides[number]["primary"],
                "engineering": guides[number]["engineering"],
                "caution": guides[number]["caution"],
                "question": guides[number]["question"],
            }
            content.update(OVERRIDES.get(number, {}))
            if number == 9:
                rebuild_slide_combined(slide, day, content)
            else:
                rebuild_slide(slide, day, content)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(args.output)
    print(f"Deck revisado: {args.output} ({len(presentation.slides)} slides; {len(REBUILD)} reconstruídos)")


if __name__ == "__main__":
    main()
